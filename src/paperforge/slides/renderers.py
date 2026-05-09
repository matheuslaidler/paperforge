"""Render typed blocks onto a python-pptx slide canvas."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from paperforge.slides.layout import estimate_para_height_emu
from paperforge.slides.parser import Block, parse_inline, parse_table
from paperforge.theme import ColorTheme


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(slide, left, top, width, height, text: str, *,
                size: int = 18, bold: bool = False,
                color: RGBColor | None = None,
                align: PP_ALIGN = PP_ALIGN.LEFT,
                font_name: str | None = None):
    color = color or RGBColor(0x22, 0x22, 0x22)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return tx


def render_para(slide, text: str, top, left, width, font_pt: int,
                theme: ColorTheme) -> int:
    h = estimate_para_height_emu(text, font_pt, int(width))
    tx = slide.shapes.add_textbox(left, top, width, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    text_color = _hex_to_rgb(theme.text)
    for chunk, bold in parse_inline(text):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(font_pt)
        r.font.bold = bold
        r.font.color.rgb = text_color
    return h


def render_quote(slide, text: str, top, left, width, font_pt: int,
                 theme: ColorTheme) -> int:
    h = estimate_para_height_emu(text, font_pt, int(width))
    tx = slide.shapes.add_textbox(left, top, width, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    muted = _hex_to_rgb(theme.muted)
    for chunk, bold in parse_inline(text):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(font_pt - 1)
        r.font.italic = True
        r.font.color.rgb = muted
    return h


def render_list(slide, items: list[str], top, left, width, font_pt: int,
                theme: ColorTheme) -> int:
    total_h = 0
    tx = slide.shapes.add_textbox(left, top, width, Inches(6))
    tf = tx.text_frame
    tf.word_wrap = True
    primary = _hex_to_rgb(theme.primary)
    text_color = _hex_to_rgb(theme.text)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        bullet = p.add_run()
        bullet.text = "• "
        bullet.font.size = Pt(font_pt)
        bullet.font.bold = True
        bullet.font.color.rgb = primary
        for chunk, bold in parse_inline(item):
            if not chunk:
                continue
            r = p.add_run()
            r.text = chunk
            r.font.size = Pt(font_pt)
            r.font.bold = bold
            r.font.color.rgb = text_color
        total_h += estimate_para_height_emu(item, font_pt, int(width))
    tx.height = max(total_h, int(Pt(font_pt * 1.5)))
    return int(tx.height)


def render_code(slide, lines: list[str], top, left, width, font_pt: int,
                theme: ColorTheme) -> int:
    n = len(lines)
    line_h = int(Pt((font_pt - 2) * 1.18))
    h = n * line_h + int(Pt(10))
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(theme.code_bg)
    bg.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    bg.shadow.inherit = False

    tx = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.05),
        width - Inches(0.2), h,
    )
    tf = tx.text_frame
    tf.word_wrap = False
    text_color = _hex_to_rgb(theme.text)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.size = Pt(font_pt - 2)
        r.font.name = "Consolas"
        r.font.color.rgb = text_color
    return h


def render_table(slide, rows: list[list[str]], top, left, width,
                 theme: ColorTheme) -> int:
    n_rows = len(rows)
    n_cols = len(rows[0])
    h = Inches(0.4 + 0.42 * n_rows)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, h).table
    primary = _hex_to_rgb(theme.primary)
    text_color = _hex_to_rgb(theme.text)
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            for chunk, bold in parse_inline(cell_text):
                run = p.add_run()
                run.text = chunk
                run.font.size = Pt(11)
                run.font.bold = bold or (i == 0)
                run.font.color.rgb = (RGBColor(0xFF, 0xFF, 0xFF)
                                      if i == 0 else text_color)
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = primary
    return h


def render_block(slide, block: Block, top: int, left, width, font_pt: int,
                 theme: ColorTheme) -> int:
    """Dispatch a block to the right renderer; returns height consumed in EMU."""
    payload = block.payload
    if block.kind == "para":
        assert isinstance(payload, str)
        return render_para(slide, payload, Emu(top), left, width, font_pt, theme)
    if block.kind == "quote":
        assert isinstance(payload, str)
        return render_quote(slide, payload, Emu(top), left, width,
                            max(font_pt - 1, 9), theme)
    if block.kind == "list":
        assert isinstance(payload, list)
        return render_list(slide, payload, Emu(top), left, width, font_pt, theme)
    if block.kind == "code":
        assert isinstance(payload, list)
        return render_code(slide, payload, Emu(top), left, width, font_pt, theme)
    if block.kind == "table":
        assert isinstance(payload, list)
        rows = parse_table(payload)
        if not rows:
            return 0
        return render_table(slide, rows, Emu(top), left, width, theme)
    return 0
