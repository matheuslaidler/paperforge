"""Cover slide construction: parse YAML frontmatter or fall back to first ``Slide 1``."""

from __future__ import annotations

import re
from dataclasses import dataclass

import yaml
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from paperforge.slides.renderers import _hex_to_rgb, add_textbox
from paperforge.theme import ColorTheme

FRONTMATTER_RE = re.compile(r"^---\n(.*?)\n---\n", re.DOTALL)


@dataclass
class CoverData:
    title: str | None = None
    subtitle: str | None = None
    author: str | None = None
    event: str | None = None
    extra: str | None = None  # arbitrary tagline


def extract_frontmatter(md_text: str) -> tuple[CoverData | None, str]:
    """Return (cover, remaining_md). If no frontmatter, returns (None, md_text)."""
    m = FRONTMATTER_RE.match(md_text)
    if not m:
        return None, md_text
    try:
        data = yaml.safe_load(m.group(1)) or {}
    except yaml.YAMLError:
        return None, md_text
    if not isinstance(data, dict):
        return None, md_text
    cover = CoverData(
        title=data.get("title"),
        subtitle=data.get("subtitle"),
        author=data.get("author"),
        event=data.get("event"),
        extra=data.get("extra"),
    )
    return cover, md_text[m.end():]


def render_cover_slide(prs, cover: CoverData, theme: ColorTheme,
                       slide_width: int, slide_height: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = _hex_to_rgb(theme.primary)
    text = _hex_to_rgb(theme.text)
    muted = _hex_to_rgb(theme.muted)

    # Vertical layout, centered.
    if cover.title:
        add_textbox(
            slide, Inches(0.6), Inches(2.3),
            slide_width - Inches(1.2), Inches(1.2),
            cover.title, size=32, bold=True, color=primary,
            align=PP_ALIGN.CENTER,
        )
    if cover.subtitle:
        add_textbox(
            slide, Inches(0.6), Inches(3.6),
            slide_width - Inches(1.2), Inches(0.9),
            cover.subtitle, size=20, color=text, align=PP_ALIGN.CENTER,
        )
    if cover.author:
        add_textbox(
            slide, Inches(0.6), slide_height - Inches(2.0),
            slide_width - Inches(1.2), Inches(0.5),
            cover.author, size=16, color=text, align=PP_ALIGN.CENTER,
        )
    if cover.event:
        add_textbox(
            slide, Inches(0.6), slide_height - Inches(1.4),
            slide_width - Inches(1.2), Inches(0.5),
            cover.event, size=13, color=muted, align=PP_ALIGN.CENTER,
        )
    if cover.extra:
        add_textbox(
            slide, Inches(0.6), slide_height - Inches(0.9),
            slide_width - Inches(1.2), Inches(0.5),
            cover.extra, size=12, color=muted, align=PP_ALIGN.CENTER,
        )
