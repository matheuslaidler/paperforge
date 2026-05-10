"""SlideBuilder — orchestrates Markdown -> PowerPoint output."""

from __future__ import annotations

import hashlib
import urllib.request
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from paperforge import i18n
from paperforge.slides.cover import (
    CoverData,
    extract_frontmatter,
    render_cover_slide,
)
from paperforge.slides.layout import fit_font_size
from paperforge.slides.parser import Block, Slide, parse_document
from paperforge.slides.renderers import (
    _hex_to_rgb,
    add_textbox,
    render_block,
)
from paperforge.theme import ColorTheme, get as get_theme
from paperforge.utils.logging import info, warn
from paperforge.utils.paths import image_cache_dir

ASPECT_RATIOS = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3": (Inches(10), Inches(7.5)),
}


@dataclass
class SlideBuildOptions:
    theme: str = "default"
    theme_file: Optional[Path] = None
    aspect: str = "16:9"
    cover: bool = True
    footer_template: str = "{n} / {total}"
    extra_browser_args: list[str] = field(default_factory=list)


@dataclass
class SlideBuildResult:
    path: Path
    slide_count: int


class SlideBuilder:
    def __init__(self, options: SlideBuildOptions | None = None) -> None:
        self.opts = options or SlideBuildOptions()
        self.theme: ColorTheme = (
            _load_theme_file(self.opts.theme_file)
            if self.opts.theme_file else get_theme(self.opts.theme)
        )

    def build(self, input_md: Path, output_pptx: Path) -> SlideBuildResult:
        if not input_md.is_file():
            raise FileNotFoundError(
                i18n.t("slides.error.input_missing", path=str(input_md))
            )
        info(i18n.t("slides.status.parsing", path=str(input_md)))
        text = input_md.read_text(encoding="utf-8")

        cover_data, remaining = extract_frontmatter(text)
        slides_data = parse_document(remaining)
        if not slides_data:
            raise ValueError(
                i18n.t("slides.error.no_slides_found", path=str(input_md))
            )

        if self.opts.aspect not in ASPECT_RATIOS:
            raise ValueError(f"Unknown aspect '{self.opts.aspect}'.")
        slide_w, slide_h = ASPECT_RATIOS[self.opts.aspect]

        prs = Presentation()
        prs.slide_width = slide_w
        prs.slide_height = slide_h

        # ---- cover ----
        if self.opts.cover:
            cover = cover_data or _synthesize_cover_from_first_slide(slides_data)
            if cover:
                render_cover_slide(prs, cover, self.theme, slide_w, slide_h)

        info(i18n.t("slides.status.rendering", n=len(slides_data)))
        total = len(slides_data) + (1 if (self.opts.cover and (cover_data or slides_data)) else 0)

        for idx, slide in enumerate(slides_data, start=1):
            self._render_slide(prs, slide, slide_w, slide_h, idx, total,
                               base_dir=input_md.parent)

        prs.save(output_pptx)
        return SlideBuildResult(path=output_pptx, slide_count=len(prs.slides))

    # ------------------------------------------------------------------ #
    #  Single-slide rendering                                            #
    # ------------------------------------------------------------------ #
    def _render_slide(self, prs, slide_data: Slide, slide_w: int, slide_h: int,
                      number: int, total: int, *, base_dir: Path) -> None:
        s = prs.slides.add_slide(prs.slide_layouts[6])

        margin = Inches(0.4)
        title_top = Inches(0.25)
        title_h = Inches(0.7)
        content_top = Inches(1.05)
        content_bottom = slide_h - Inches(0.5)
        footer_top = slide_h - Inches(0.45)

        primary = _hex_to_rgb(self.theme.primary)
        muted = _hex_to_rgb(self.theme.muted)

        # title
        add_textbox(
            s, margin, title_top, slide_w - 2 * margin, title_h,
            slide_data.title, size=24, bold=True, color=primary,
        )

        text_blocks: list[Block] = []
        image_blocks: list[Block] = []
        for b in slide_data.blocks:
            (image_blocks if b.kind == "image" else text_blocks).append(b)

        has_img = bool(image_blocks)
        text_left = margin
        text_width = Inches(7.0) if has_img else (slide_w - 2 * margin)
        available_h = content_bottom - content_top

        font_pt = fit_font_size(text_blocks, int(text_width), int(available_h))
        if font_pt < 15 and text_blocks:
            warn(i18n.t("slides.warn.shrink", n=number, pt=font_pt))

        cur = int(content_top)
        bottom = int(content_bottom)
        gap = int(Inches(0.08))
        for block in text_blocks:
            if cur >= bottom:
                break
            try:
                used = render_block(s, block, cur, text_left, text_width,
                                    font_pt, self.theme)
            except Exception as exc:  # noqa: BLE001
                warn(f"Block ({block.kind}) failed: {exc}")
                used = 0
            cur += int(used) + gap

        # right column: first image (resolves URL via cache).
        if has_img:
            ref = str(image_blocks[0].payload)
            img_path = self._resolve_image(ref, base_dir)
            if img_path:
                s.shapes.add_picture(
                    str(img_path), Inches(7.7), Inches(1.1),
                    width=slide_w - Inches(7.9),
                )

        # footer
        footer_text = self.opts.footer_template.format(n=number, total=total)
        add_textbox(
            s, slide_w - Inches(1.6), footer_top, Inches(1.4), Inches(0.3),
            footer_text, size=10, color=muted, align=PP_ALIGN.RIGHT,
        )

    # ------------------------------------------------------------------ #
    #  Image resolution (local + URL cache)                              #
    # ------------------------------------------------------------------ #
    def _resolve_image(self, reference: str, base_dir: Path) -> Path | None:
        if reference.startswith(("http://", "https://")):
            cached = self._url_cache_path(reference)
            if cached.is_file():
                return cached
            try:
                urllib.request.urlretrieve(reference, cached)  # noqa: S310
                return cached if cached.is_file() else None
            except Exception:
                warn(f"Could not fetch image: {reference}")
                return None
        # local
        candidate = (base_dir / reference).resolve()
        return candidate if candidate.is_file() else None

    @staticmethod
    def _url_cache_path(url: str) -> Path:
        digest = hashlib.sha1(url.encode("utf-8")).hexdigest()[:16]
        suffix = Path(url.split("?")[0]).suffix or ".img"
        return image_cache_dir() / f"{digest}{suffix}"


# --------------------------------------------------------------------------- #
#  Helpers                                                                    #
# --------------------------------------------------------------------------- #
def _load_theme_file(path: Path) -> ColorTheme:
    import yaml
    with path.open(encoding="utf-8") as fh:
        data = yaml.safe_load(fh) or {}
    return ColorTheme(**{k: v for k, v in data.items()
                         if k in ColorTheme.__dataclass_fields__})


def _synthesize_cover_from_first_slide(slides: list[Slide]) -> CoverData | None:
    if not slides:
        return None
    first = slides[0]
    title = first.title
    subtitle = None
    for b in first.blocks:
        if b.kind == "para" and isinstance(b.payload, str) and b.payload.strip():
            subtitle = b.payload.strip()
            break
    return CoverData(title=title, subtitle=subtitle)
